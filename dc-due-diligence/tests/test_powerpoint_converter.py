"""Tests for the PowerPoint presentation converter."""

from pathlib import Path

import pptx
import pytest
from pptx.util import Inches

from converters import ConfidenceLevel, ExtractionResult
from converters.powerpoint import (
    PowerPointConverter,
    _clean_cell,
    _score_confidence,
    _text_frame_to_markdown,
)


# ------------------------------------------------------------------
# Fixtures
# ------------------------------------------------------------------


@pytest.fixture
def converter():
    return PowerPointConverter()


@pytest.fixture
def simple_pptx(tmp_path):
    """Create a simple .pptx with a title slide and a content slide."""
    prs = pptx.Presentation()
    layout = prs.slide_layouts[1]  # Title and Content layout

    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Data Center Overview"
    slide1.placeholders[1].text = "Pioneer Park facility in Dallas, TX."

    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Power Details"
    slide2.placeholders[1].text = "50MW available with Oncor service.\nExpansion to 130MW planned."

    path = tmp_path / "simple.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def pptx_with_table(tmp_path):
    """Create a .pptx with a slide containing a table."""
    prs = pptx.Presentation()
    layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(layout)

    # Add a title manually.
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    txBox.text_frame.paragraphs[0].text = "Site Specifications"

    # Add a table.
    rows, cols = 4, 2
    table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(6), Inches(3))
    table = table_shape.table
    table.cell(0, 0).text = "Attribute"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Total Power"
    table.cell(1, 1).text = "50 MW"
    table.cell(2, 0).text = "Floor Space"
    table.cell(2, 1).text = "120,000 sqft"
    table.cell(3, 0).text = "PUE"
    table.cell(3, 1).text = "1.25"

    path = tmp_path / "with_table.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def empty_pptx(tmp_path):
    """Create a .pptx with no slides."""
    prs = pptx.Presentation()
    path = tmp_path / "empty.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def pptx_blank_slides(tmp_path):
    """Create a .pptx where slides have titles but no body content."""
    prs = pptx.Presentation()
    layout = prs.slide_layouts[6]  # Blank layout
    prs.slides.add_slide(layout)
    prs.slides.add_slide(layout)
    path = tmp_path / "blank_slides.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def pptx_with_image(tmp_path):
    """Create a .pptx with an embedded image."""
    from PIL import Image as PILImage
    import io

    # Create a small test image and save to file.
    img = PILImage.new("RGB", (100, 100), color="blue")
    img_path = tmp_path / "test_image.png"
    img.save(str(img_path))

    prs = pptx.Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Site Photos"
    slide.placeholders[1].text = "Front view of the facility."
    slide.shapes.add_picture(str(img_path), Inches(1), Inches(3), Inches(4), Inches(3))

    path = tmp_path / "with_image.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def pptx_with_notes(tmp_path):
    """Create a .pptx with speaker notes."""
    prs = pptx.Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Key Points"
    slide.placeholders[1].text = "Revenue growth projections."

    # Add speaker notes.
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Mention the Q3 numbers and Oncor timeline."

    path = tmp_path / "with_notes.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def pptx_many_slides(tmp_path):
    """Create a .pptx with many slides to test confidence scoring."""
    prs = pptx.Presentation()
    layout = prs.slide_layouts[1]
    for i in range(10):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"Section {i + 1}"
        slide.placeholders[1].text = f"Content for section {i + 1} with details about the data center facility."

    path = tmp_path / "many_slides.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def pptx_special_chars(tmp_path):
    """Create a .pptx with a table containing pipe chars."""
    prs = pptx.Presentation()
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    txBox.text_frame.paragraphs[0].text = "Special Characters"

    rows, cols = 2, 2
    table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(6), Inches(2))
    table = table_shape.table
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Test|Item"
    table.cell(1, 1).text = "A\nB"

    path = tmp_path / "special.pptx"
    prs.save(str(path))
    return path


# ------------------------------------------------------------------
# Unit tests for helper functions
# ------------------------------------------------------------------


class TestCleanCell:
    def test_strips_whitespace(self):
        assert _clean_cell("  padded  ") == "padded"

    def test_replaces_newlines(self):
        assert _clean_cell("line1\nline2") == "line1 line2"

    def test_replaces_carriage_return(self):
        assert _clean_cell("line1\rline2") == "line1 line2"

    def test_empty_string(self):
        assert _clean_cell("") == ""


class TestScoreConfidence:
    def test_no_slides(self):
        confidence, reason = _score_confidence(0, 0, 0)
        assert confidence == ConfidenceLevel.LOW
        assert "no slides" in reason

    def test_no_content(self):
        confidence, reason = _score_confidence(0, 5, 0)
        assert confidence == ConfidenceLevel.LOW
        assert "no text content" in reason

    def test_very_little_text(self):
        confidence, reason = _score_confidence(10, 3, 2)
        assert confidence == ConfidenceLevel.LOW
        assert "very little" in reason

    def test_sparse_content(self):
        confidence, reason = _score_confidence(80, 10, 2)
        assert confidence == ConfidenceLevel.MEDIUM
        assert "sparse" in reason

    def test_substantial_content(self):
        confidence, reason = _score_confidence(500, 5, 5)
        assert confidence == ConfidenceLevel.HIGH

    def test_low_content_ratio(self):
        confidence, reason = _score_confidence(200, 20, 3)
        assert confidence == ConfidenceLevel.MEDIUM


# ------------------------------------------------------------------
# Converter behaviour tests
# ------------------------------------------------------------------


class TestPowerPointConverterBasics:
    def test_supported_extensions(self, converter):
        assert ".pptx" in converter.supported_extensions

    def test_can_handle_pptx(self, converter):
        assert converter.can_handle(Path("presentation.pptx")) is True
        assert converter.can_handle(Path("PRESENTATION.PPTX")) is True

    def test_cannot_handle_non_pptx(self, converter):
        assert converter.can_handle(Path("data.pdf")) is False
        assert converter.can_handle(Path("data.xlsx")) is False
        assert converter.can_handle(Path("data.docx")) is False
        assert converter.can_handle(Path("data.ppt")) is False

    def test_missing_file_returns_error(self, converter):
        result = converter.convert(Path("/tmp/nonexistent_pptx_12345.pptx"))
        assert result.success is False
        assert result.confidence == ConfidenceLevel.LOW
        assert "not found" in result.error.lower()
        assert result.method == "python-pptx"

    def test_wrong_extension_returns_error(self, converter, tmp_path):
        txt = tmp_path / "test.txt"
        txt.write_text("not a pptx")
        result = converter.convert(txt)
        assert result.success is False
        assert "not a powerpoint" in result.error.lower()

    def test_result_type(self, converter):
        result = converter.convert(Path("/tmp/nonexistent_12345.pptx"))
        assert isinstance(result, ExtractionResult)

    def test_source_path_resolved(self, converter, simple_pptx):
        result = converter.convert(simple_pptx)
        assert result.source_path == simple_pptx.resolve()


class TestPptxExtraction:
    def test_simple_file(self, converter, simple_pptx):
        result = converter.convert(simple_pptx)
        assert result.success is True
        assert result.method == "python-pptx"

    def test_contains_slide_titles(self, converter, simple_pptx):
        result = converter.convert(simple_pptx)
        assert "Data Center Overview" in result.text
        assert "Power Details" in result.text

    def test_slide_numbering(self, converter, simple_pptx):
        result = converter.convert(simple_pptx)
        assert "Slide 1" in result.text
        assert "Slide 2" in result.text

    def test_contains_body_text(self, converter, simple_pptx):
        result = converter.convert(simple_pptx)
        assert "Pioneer Park" in result.text
        assert "50MW" in result.text

    def test_slide_separator(self, converter, simple_pptx):
        result = converter.convert(simple_pptx)
        # Slides should be separated by horizontal rules.
        assert "---" in result.text

    def test_table_extraction(self, converter, pptx_with_table):
        result = converter.convert(pptx_with_table)
        assert result.success is True
        assert "Attribute" in result.text
        assert "50 MW" in result.text
        assert "120,000 sqft" in result.text
        assert "PUE" in result.text
        assert "1.25" in result.text

    def test_table_is_markdown_formatted(self, converter, pptx_with_table):
        result = converter.convert(pptx_with_table)
        assert "|" in result.text
        assert "---" in result.text

    def test_page_count_matches_slides(self, converter, simple_pptx):
        result = converter.convert(simple_pptx)
        assert result.page_count == 2

    def test_metadata_present(self, converter, simple_pptx):
        result = converter.convert(simple_pptx)
        assert "total_slides" in result.metadata
        assert "slides_with_content" in result.metadata
        assert "slides_empty" in result.metadata
        assert "total_text_shapes" in result.metadata
        assert "total_tables" in result.metadata
        assert "total_images" in result.metadata
        assert "total_charts" in result.metadata
        assert "total_chars" in result.metadata

    def test_metadata_values(self, converter, simple_pptx):
        result = converter.convert(simple_pptx)
        assert result.metadata["total_slides"] == 2
        assert result.metadata["slides_with_content"] == 2
        assert result.metadata["slides_empty"] == 0


class TestPptxImages:
    def test_image_noted(self, converter, pptx_with_image):
        result = converter.convert(pptx_with_image)
        assert result.success is True
        assert "[Image]" in result.text
        assert result.metadata["total_images"] >= 1

    def test_no_image_note_when_none(self, converter, simple_pptx):
        result = converter.convert(simple_pptx)
        assert "[Image]" not in result.text
        assert result.metadata["total_images"] == 0


class TestPptxNotes:
    def test_speaker_notes_extracted(self, converter, pptx_with_notes):
        result = converter.convert(pptx_with_notes)
        assert result.success is True
        assert "Speaker Notes" in result.text
        assert "Q3 numbers" in result.text

    def test_no_notes_when_absent(self, converter, simple_pptx):
        result = converter.convert(simple_pptx)
        assert "Speaker Notes" not in result.text


class TestPptxEdgeCases:
    def test_empty_presentation(self, converter, empty_pptx):
        result = converter.convert(empty_pptx)
        assert result.success is True
        assert result.confidence == ConfidenceLevel.LOW
        assert result.page_count == 0

    def test_blank_slides(self, converter, pptx_blank_slides):
        result = converter.convert(pptx_blank_slides)
        assert result.success is True
        # Blank slides should not count as having content.
        assert result.metadata["slides_with_content"] == 0

    def test_many_slides(self, converter, pptx_many_slides):
        result = converter.convert(pptx_many_slides)
        assert result.success is True
        assert result.confidence == ConfidenceLevel.HIGH
        assert result.page_count == 10
        assert result.metadata["slides_with_content"] == 10

    def test_special_characters_in_table(self, converter, pptx_special_chars):
        result = converter.convert(pptx_special_chars)
        assert result.success is True
        assert "Test\\|Item" in result.text
        assert "A B" in result.text

    def test_corrupted_file_handled(self, converter, tmp_path):
        bad_file = tmp_path / "corrupted.pptx"
        bad_file.write_bytes(b"this is not a real pptx file")
        result = converter.convert(bad_file)
        assert result.success is False
        assert result.confidence == ConfidenceLevel.LOW
        assert result.error is not None

    def test_confidence_summary_works(self, converter, simple_pptx):
        result = converter.convert(simple_pptx)
        summary = result.confidence_summary
        assert "python-pptx" in summary
        assert "confidence" in summary


class TestPptxImportFromPackage:
    def test_import_from_converters(self):
        from converters import PowerPointConverter
        converter = PowerPointConverter()
        assert ".pptx" in converter.supported_extensions
