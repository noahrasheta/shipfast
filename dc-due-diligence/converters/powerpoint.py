"""
PowerPoint presentation conversion for .pptx files.

Extracts slide titles and text content from each slide using python-pptx.
Output is clean markdown with each slide as a numbered section, slide
titles as headings, and text content as paragraphs.  Notes the presence
of embedded images and charts.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.presentation import Presentation
from pptx.slide import Slide

from converters.base import BaseConverter, ConfidenceLevel, ExtractionResult


class PowerPointConverter(BaseConverter):
    """Extract text from PowerPoint presentations (.pptx) using python-pptx.

    Each slide is rendered as a markdown section with a heading for the
    slide title and paragraphs for the body text.  Tables within slides
    are rendered as markdown tables.  Embedded images and charts are
    noted with placeholders.
    """

    supported_extensions: list[str] = [".pptx"]

    def convert(self, path: Path) -> ExtractionResult:
        """Open a .pptx file, extract all slides, and return markdown."""
        path = Path(path).resolve()

        if not path.exists():
            return ExtractionResult(
                source_path=path,
                text="",
                method="python-pptx",
                success=False,
                confidence=ConfidenceLevel.LOW,
                confidence_reason="file not found",
                error=f"File not found: {path}",
            )

        if path.suffix.lower() not in self.supported_extensions:
            return ExtractionResult(
                source_path=path,
                text="",
                method="python-pptx",
                success=False,
                confidence=ConfidenceLevel.LOW,
                confidence_reason="unsupported file type",
                error=f"Not a PowerPoint file: {path.suffix}",
            )

        try:
            return self._extract(path)
        except Exception as exc:
            return ExtractionResult(
                source_path=path,
                text="",
                method="python-pptx",
                success=False,
                confidence=ConfidenceLevel.LOW,
                confidence_reason="extraction crashed",
                error=f"Extraction failed: {exc}",
            )

    # ------------------------------------------------------------------
    # Internal extraction
    # ------------------------------------------------------------------

    def _extract(self, path: Path) -> ExtractionResult:
        """Core extraction -- iterate slides and build markdown."""
        presentation = pptx.Presentation(str(path))

        slide_sections: list[str] = []
        total_chars = 0
        total_slides = len(presentation.slides)
        slides_with_content = 0
        total_text_shapes = 0
        total_table_count = 0
        total_image_count = 0
        total_chart_count = 0

        for slide_num, slide in enumerate(presentation.slides, start=1):
            slide_md, stats = _extract_slide(slide, slide_num)
            if slide_md:
                slide_sections.append(slide_md)
                total_chars += len(slide_md)
                slides_with_content += 1
                total_text_shapes += stats["text_shapes"]
                total_table_count += stats["tables"]
                total_image_count += stats["images"]
                total_chart_count += stats["charts"]

        combined_text = "\n\n---\n\n".join(slide_sections)

        # Confidence scoring.
        confidence, confidence_reason = _score_confidence(
            total_chars=total_chars,
            total_slides=total_slides,
            slides_with_content=slides_with_content,
        )

        metadata: dict[str, Any] = {
            "total_slides": total_slides,
            "slides_with_content": slides_with_content,
            "slides_empty": total_slides - slides_with_content,
            "total_text_shapes": total_text_shapes,
            "total_tables": total_table_count,
            "total_images": total_image_count,
            "total_charts": total_chart_count,
            "total_chars": total_chars,
        }

        return ExtractionResult(
            source_path=path,
            text=combined_text,
            method="python-pptx",
            success=True,
            confidence=confidence,
            confidence_reason=confidence_reason,
            page_count=total_slides,
            metadata=metadata,
        )


# ------------------------------------------------------------------
# Module-level helpers
# ------------------------------------------------------------------


def _extract_slide(
    slide: Slide, slide_num: int
) -> tuple[str, dict[str, int]]:
    """Extract all content from a single slide.

    Returns (markdown_text, stats_dict).  The stats dict contains counts
    of text_shapes, tables, images, and charts found on this slide.
    """
    parts: list[str] = []
    title_text = ""
    stats = {"text_shapes": 0, "tables": 0, "images": 0, "charts": 0}

    # Try to get the slide title from the title placeholder.
    if slide.shapes.title is not None:
        title_text = slide.shapes.title.text.strip()

    # Build slide heading.
    if title_text:
        parts.append(f"## Slide {slide_num}: {title_text}")
    else:
        parts.append(f"## Slide {slide_num}")

    # Process all shapes on the slide.
    for shape in slide.shapes:
        # Skip the title shape -- we already extracted it above.
        if shape.has_text_frame and shape == slide.shapes.title:
            stats["text_shapes"] += 1
            continue

        if shape.has_table:
            md_table = _shape_table_to_markdown(shape.table)
            if md_table:
                parts.append(md_table)
                stats["tables"] += 1
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            stats["charts"] += 1
            parts.append("*[Chart]*")
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            stats["images"] += 1
            parts.append("*[Image]*")
            continue

        # For group shapes, check for images inside.
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_images = _count_group_images(shape)
            if group_images > 0:
                stats["images"] += group_images
                parts.append(
                    f"*[{group_images} grouped image"
                    f"{'s' if group_images != 1 else ''}]*"
                )

        if shape.has_text_frame:
            text_md = _text_frame_to_markdown(shape.text_frame)
            if text_md:
                parts.append(text_md)
                stats["text_shapes"] += 1

    # Check for speaker notes.
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            parts.append(f"\n**Speaker Notes:** {notes_text}")

    # If the slide only has the heading and nothing else, it's empty.
    if len(parts) <= 1:
        return "", stats

    return "\n\n".join(parts), stats


def _text_frame_to_markdown(text_frame: Any) -> str:
    """Extract text from a shape's text frame, preserving paragraph breaks."""
    paragraphs: list[str] = []
    for para in text_frame.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)

    return "\n\n".join(paragraphs)


def _shape_table_to_markdown(table: Any) -> str:
    """Convert a PowerPoint table shape to a markdown table string."""
    rows: list[list[str]] = []
    for row in table.rows:
        cells = [_clean_cell(cell.text) for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ""

    # Skip entirely empty tables.
    if not any(cell for row in rows for cell in row):
        return ""

    max_cols = max(len(r) for r in rows)
    if max_cols == 0:
        return ""

    # Normalize column count.
    normalized: list[list[str]] = []
    for row in rows:
        padded = list(row)
        while len(padded) < max_cols:
            padded.append("")
        normalized.append(padded[:max_cols])

    # Escape pipes.
    escaped: list[list[str]] = []
    for row in normalized:
        escaped.append([cell.replace("|", "\\|") for cell in row])

    header = escaped[0]
    separator = ["---"] * max_cols
    body = escaped[1:]

    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    for row in body:
        lines.append("| " + " | ".join(row) + " |")

    return "\n".join(lines)


def _clean_cell(text: str) -> str:
    """Clean a table cell's text for markdown display."""
    cleaned = text.strip()
    cleaned = cleaned.replace("\n", " ").replace("\r", " ")
    return cleaned


def _count_group_images(shape: Any) -> int:
    """Count picture shapes inside a group shape."""
    count = 0
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                count += 1
            elif child.shape_type == MSO_SHAPE_TYPE.GROUP:
                count += _count_group_images(child)
    return count


def _score_confidence(
    total_chars: int,
    total_slides: int,
    slides_with_content: int,
) -> tuple[ConfidenceLevel, str]:
    """Determine confidence based on extraction metrics.

    Returns (confidence, reason).
    """
    if total_slides == 0:
        return (
            ConfidenceLevel.LOW,
            "presentation contains no slides",
        )

    if slides_with_content == 0:
        return (
            ConfidenceLevel.LOW,
            f"no text content found across {total_slides} slide(s)",
        )

    if total_chars < 20:
        return (
            ConfidenceLevel.LOW,
            f"very little text, only {total_chars} characters "
            f"from {slides_with_content} slide(s)",
        )

    content_ratio = slides_with_content / total_slides
    if total_chars < 100 or content_ratio < 0.3:
        return (
            ConfidenceLevel.MEDIUM,
            f"sparse content, {slides_with_content} of "
            f"{total_slides} slide(s) had text",
        )

    return (
        ConfidenceLevel.HIGH,
        f"{slides_with_content} of {total_slides} slide(s) "
        f"extracted with {total_chars} total characters",
    )
